import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.util import Inches


# =========================
# 配置区：你只改这里
# =========================
DECKTAPE_CMD = r"D:\Node\node_cache\decktape.cmd"

URLS_FILE = Path(r"D:\Desktop\Neural Symbolic\Code\20260106\urls.txt")
OUT_ROOT = Path(r"D:\Desktop\Neural Symbolic\Code\20260106\batch_output").resolve()

# 16:9 高清（你想要更高也行，比如 2560x1440 / 3200x1800）
SIZE = "2560x1440"

FRAGMENTS = False          # True: 导出逐步出现（页数变多）
LOAD_PAUSE_MS = 800
PAUSE_MS = 800

# 网络慢/MathJax 多时建议加大
URL_LOAD_TIMEOUT_MS = 120000
PAGE_LOAD_TIMEOUT_MS = 60000
BUFFER_TIMEOUT_MS = 60000

# 只跑前 N 个用于测试：None 表示全跑
LIMIT = None
# =========================


def natural_sort_key(s: str):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", s)]


def safe_stem_from_url(url: str, max_len: int = 80) -> str:
    u = urlparse(url)
    host = u.netloc.replace(":", "_")
    path = (u.path or "").strip("/").replace("/", "_")
    frag = (u.fragment or "").replace("/", "_").replace("\\", "_")
    base = "_".join([p for p in [host, path, frag] if p])

    base = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", base)
    base = re.sub(r"_+", "_", base).strip("._-")
    if not base:
        base = "slides"
    return base[:max_len]


def read_urls(file_path: Path):
    if not file_path.exists():
        raise FileNotFoundError(f"找不到 urls 文件：{file_path}")
    urls = []
    for line in file_path.read_text(encoding="utf-8").splitlines():
        s = line.strip()
        if not s or s.startswith("#"):
            continue
        urls.append(s)
    return urls


def make_tmp_base() -> Path:
    # 临时目录也放在 OUT_ROOT 所在盘，避免跨盘移动
    tmp_base = OUT_ROOT / "_tmp"
    tmp_base.mkdir(parents=True, exist_ok=True)
    return tmp_base


def decktape_export_png(url: str, png_dir: Path):
    png_dir.mkdir(parents=True, exist_ok=True)
    tmp_base = make_tmp_base()

    with tempfile.TemporaryDirectory(dir=str(tmp_base)) as td:
        td_path = Path(td)

        # 关键：filename 参数只给纯文件名，避免 Windows 盘符污染截图路径
        pdf_name = "out.pdf"

        cmd = [
            DECKTAPE_CMD,
            "reveal",
            "--screenshots",
            "--screenshots-directory",
            str(png_dir),
            "--screenshots-size",
            SIZE,
            "--screenshots-format",
            "png",
            "--pause",
            str(PAUSE_MS),
            "--load-pause",
            str(LOAD_PAUSE_MS),
            "--url-load-timeout",
            str(URL_LOAD_TIMEOUT_MS),
            "--page-load-timeout",
            str(PAGE_LOAD_TIMEOUT_MS),
            "--buffer-timeout",
            str(BUFFER_TIMEOUT_MS),
        ]
        if FRAGMENTS:
            cmd.append("--fragments")

        cmd += [url, pdf_name]

        print(f"[INFO] PNG 导出: {png_dir}")
        subprocess.run(cmd, check=True, cwd=str(td_path))


def decktape_export_pdf(url: str, out_pdf: Path):
    out_pdf.parent.mkdir(parents=True, exist_ok=True)
    tmp_base = make_tmp_base()

    with tempfile.TemporaryDirectory(dir=str(tmp_base)) as td:
        td_path = Path(td)

        # 同样：只给纯文件名
        pdf_name = out_pdf.name

        cmd = [
            DECKTAPE_CMD,
            "reveal",
            "--size",
            SIZE,
            "--pause",
            str(PAUSE_MS),
            "--load-pause",
            str(LOAD_PAUSE_MS),
            "--url-load-timeout",
            str(URL_LOAD_TIMEOUT_MS),
            "--page-load-timeout",
            str(PAGE_LOAD_TIMEOUT_MS),
            "--buffer-timeout",
            str(BUFFER_TIMEOUT_MS),
        ]
        if FRAGMENTS:
            cmd.append("--fragments")

        cmd += [url, pdf_name]

        print(f"[INFO] PDF 导出: {out_pdf}")
        subprocess.run(cmd, check=True, cwd=str(td_path))

        generated = td_path / pdf_name
        if not generated.exists():
            raise FileNotFoundError(f"未找到生成的 PDF：{generated}")

        # 用 move，稳（即使未来你改成跨盘也不怕）
        shutil.move(str(generated), str(out_pdf))


def build_pptx_from_pngs(png_dir: Path, pptx_path: Path):
    pngs = sorted(
        [p for p in png_dir.iterdir() if p.suffix.lower() == ".png"],
        key=lambda p: natural_sort_key(p.name),
    )
    if not pngs:
        raise RuntimeError(f"在目录里没找到 PNG：{png_dir}")

    prs = Presentation()
    # 16:9 页面（与你的目标一致）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for img_path in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img_path), 0, 0, width=slide_w, height=slide_h)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"[INFO] PPTX 已生成: {pptx_path}")


def main():
    if not Path(DECKTAPE_CMD).exists():
        print(f"[ERROR] decktape.cmd 不存在：{DECKTAPE_CMD}")
        sys.exit(1)

    urls = read_urls(URLS_FILE)
    if LIMIT is not None:
        urls = urls[:LIMIT]

    OUT_ROOT.mkdir(parents=True, exist_ok=True)

    ok, fail = 0, 0
    for i, url in enumerate(urls, start=1):
        stem = f"{i:03d}_" + safe_stem_from_url(url)
        one_dir = OUT_ROOT / stem

        png_dir = one_dir / "png"
        pdf_path = one_dir / f"{stem}.pdf"
        pptx_path = one_dir / f"{stem}.pptx"

        print(f"\n===== [{i}/{len(urls)}] {url} =====")
        try:
            # 1) PNG
            decktape_export_png(url, png_dir)

            # 2) PPTX
            build_pptx_from_pngs(png_dir, pptx_path)

            # 3) PDF
            decktape_export_pdf(url, pdf_path)

            ok += 1
        except subprocess.CalledProcessError as e:
            fail += 1
            print(f"[ERROR] decktape 失败: {url}\n{e}\n")
        except Exception as e:
            fail += 1
            print(f"[ERROR] 失败: {url}\n{e}\n")

    print(f"\n[DONE] 成功 {ok} 个，失败 {fail} 个。输出目录：{OUT_ROOT}")


if __name__ == "__main__":
    main()
