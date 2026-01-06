# reveal_to_pptx.py
# Reveal.js 网页课件 -> PNG -> PPTX（每页一张图，版式一致）

import re
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


# =========================
# 配置区：你只改这里就行
# =========================
DECKTAPE_CMD = r"D:\Node\node_cache\decktape.cmd"
URL = "https://daiwz.net/course/symb_learn/2025/0-1.html#/slide-title"

OUT_DIR = Path(r"D:\Desktop\Neural Symbolic\Code\20260106\output").resolve()
SIZE = "1920x1080"
FRAGMENTS = False

LOAD_PAUSE_MS = 800
PAUSE_MS = 800
# =========================


def natural_sort_key(s: str):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", s)]


def run_decktape_screenshots(
    url: str,
    out_dir: Path,
    size: str,
    fragments: bool,
    pause_ms: int,
    load_pause_ms: int,
):
    out_dir.mkdir(parents=True, exist_ok=True)

    decktape_path = Path(DECKTAPE_CMD)
    if not decktape_path.exists():
        raise FileNotFoundError(f"找不到 decktape.cmd：{DECKTAPE_CMD}")

    # 关键修复点：
    # 1) decktape 的 <filename> 参数不要给绝对路径（带 C:\ 这种会被拼进截图名，导致非法路径）
    # 2) 用 cwd=临时目录，让 out.pdf 落在临时目录里
    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        pdf_name = "out.pdf"  # 只给纯文件名，避免 Windows 盘符污染截图名

        cmd = [
            DECKTAPE_CMD,
            "reveal",
            "--screenshots",
            "--screenshots-directory",
            str(out_dir),
            "--screenshots-size",
            size,
            "--screenshots-format",
            "png",
            "--pause",
            str(pause_ms),
            "--load-pause",
            str(load_pause_ms),
        ]
        if fragments:
            cmd.append("--fragments")

        cmd += [url, pdf_name]

        print("[INFO] 运行 Decktape 截图中...")
        print("[INFO] 命令：", " ".join(cmd))
        subprocess.run(cmd, check=True, cwd=str(td_path))
        print("[INFO] 截图完成：", out_dir)


def build_pptx_from_pngs(png_dir: Path, pptx_path: Path, widescreen_16_9: bool = True):
    pngs = sorted(
        [p for p in png_dir.iterdir() if p.suffix.lower() == ".png"],
        key=lambda p: natural_sort_key(p.name),
    )
    if not pngs:
        raise RuntimeError(f"在目录里没找到 PNG：{png_dir}")

    prs = Presentation()
    if widescreen_16_9:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    print(f"[INFO] 生成 PPT：共 {len(pngs)} 页")
    for i, img_path in enumerate(pngs, start=1):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img_path), 0, 0, width=slide_w, height=slide_h)
        if i % 10 == 0:
            print(f"[INFO] 已写入 {i}/{len(pngs)} 页")

    prs.save(str(pptx_path))
    print("[INFO] PPT 已保存：", pptx_path)


def main():
    png_dir = OUT_DIR / "png"
    pptx_path = OUT_DIR / "slides.pptx"
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    run_decktape_screenshots(
        url=URL,
        out_dir=png_dir,
        size=SIZE,
        fragments=FRAGMENTS,
        pause_ms=PAUSE_MS,
        load_pause_ms=LOAD_PAUSE_MS,
    )

    build_pptx_from_pngs(png_dir, pptx_path)


if __name__ == "__main__":
    try:
        main()
    except subprocess.CalledProcessError as e:
        print("[ERROR] Decktape 运行失败：", e)
        sys.exit(2)
    except Exception as e:
        print("[ERROR] 失败：", e)
        sys.exit(3)
